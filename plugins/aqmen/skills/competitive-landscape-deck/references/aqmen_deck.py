"""aqmen_deck — build decision-grade CDD slide decks in aqmen's house style.

A thin, opinionated wrapper around python-pptx that encodes aqmen's deck design
system (the look of the CDD "Discussion Materials" deliverable): 16:9 slides,
Montserrat on a blue-dominant navy palette, and the standard slide chrome (DRAFT tag,
section eyebrow, wordmark, page number, source line). It exposes a small set of
slide builders — title, agenda/divider, executive summary, content-with-takeaways,
and native charts (column / bar / line / stacked, plus a variable-width Marimekko
and harvey-ball matrix) — so every deck comes out consistent and editable.

This module is CANONICAL in `shared/` and synced into each deck skill's
`references/`. Edit it here and re-run `scripts/sync-shared.mjs`.

Requires: python-pptx (`pip install python-pptx`). Charts are NATIVE PowerPoint
charts (editable in PowerPoint / Google Slides), not images.

Minimal usage:

    from aqmen_deck import Deck, Bullet

    d = Deck(draft=True)
    d.title_slide("The CDD Market", "A Strategic Opportunity", "2025")
    d.agenda(["Context", "Market Overview", "Competitive Landscape", "Appendix"],
             active="Market Overview")
    d.content_slide(
        eyebrow=("Market Overview", "TAM"),
        headline="The CDD TAM is a ~$90B opportunity led by North America",
        body=[Bullet("North America is ~two-thirds of total value", level=0, bold=True)],
        takeaways=["NA combines high deal volume and premium pricing",
                   "Europe is second-largest despite fewer deals"],
        source="GainPro, Bloomberg, aqmen", illustrative=True)
    d.save("out.pptx")
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from typing import Iterable, Sequence

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN_TYPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as _RT
from pptx.oxml.ns import qn as _qn
from pptx.shapes.autoshape import Shape as _Shape
from pptx.util import Emu, Inches, Pt

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# --------------------------------------------------------------------------- #
# Design tokens — extracted from the aqmen brand theme + CDD deliverable deck. #
# --------------------------------------------------------------------------- #


def _c(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


class Palette:
    """aqmen brand palette (hex without '#'). Blue-dominant: navy leads, with a
    coherent ramp of blues; a vivid azure is the single accent."""

    NAVY = _c("03045E")        # primary brand navy — titles, wordmark, headlines
    NAVY_DEEP = _c("020462")   # deepest navy — exec-summary band, dividers
    INK = _c("302E2E")         # body text
    BLUE = _c("0728A3")        # secondary blue — data series, subtitles, links
    AZURE = _c("2E6FD6")       # vivid azure — the accent (DRAFT tag, cover bar)
    STEEL = _c("6BAED6")       # light steel blue — mid chart series
    CYAN = _c("ADE8F3")        # light cyan — fills, secondary bands (e.g. SAM)
    CYAN_MUTED = _c("D6EEF5")  # very light cyan — highlight rows, agenda active
    GREEN = _c("00FF7F")       # spring green — use sparingly, positive deltas
    WHITE = _c("FEFFFF")
    GREY = _c("8A8A8A")        # source lines, page numbers, captions
    GREY_LINE = _c("C9CBD6")   # hairlines / dividers
    GREY_BG = _c("F2F3F7")     # panel backgrounds
    AMBER = _c("B7791F")       # medium-certainty dot / caution
    RED = _c("9B1C1C")         # low-certainty dot
    ORANGE = _c("E4572E")      # from→to comparison, negative bridge steps

    # Certainty dots for driver trees: 0 low → 1 medium → 2 high.
    CERTAINTY = [RED, AMBER, NAVY]

    # Ordered categorical series for charts — a navy-led blue ramp, grey neutral.
    SERIES = [NAVY, BLUE, AZURE, STEEL, CYAN, GREY]


class Font:
    HEAD = "Montserrat"          # headlines / titles (rendered SemiBold/Bold)
    BODY = "Montserrat Medium"   # body copy
    FALLBACK = "Arial"           # portable fallback if Montserrat is absent
    HEAD_THEME = "Montserrat"    # theme major font (baked into the template)
    BODY_THEME = "Montserrat"    # theme minor font


# Slide geometry (16:9, 13.333in x 7.5in — matches the CDD deck page size).
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.95)      # first y below the headline band
CONTENT_BOTTOM = Inches(6.95)   # y above the footer
RAIL_X = Inches(9.05)           # left edge of the right "Key Takeaways" rail
RAIL_W = Inches(3.75)


# --------------------------------------------------------------------------- #
# Branding — applied in code so a blank deck comes out fully on-brand, with no  #
# dependency on an external base file. (python-pptx always ships lxml.)         #
# --------------------------------------------------------------------------- #

_THEME_CLR = {
    "dk1": "020462", "lt1": "FEFFFF", "dk2": "302E2E", "lt2": "ADE8F3",
    "accent1": "03045E", "accent2": "ADE8F3", "accent3": "0728A3",
    "accent4": "00FF7F", "accent5": "2E6FD6", "accent6": "302E2E",
    "hlink": "00FF7F", "folHlink": "0728A3",
}


def _clr_scheme():
    ns = f"{{{_A}}}"
    scheme = etree.Element(f"{ns}clrScheme")
    scheme.set("name", "Aqmen")
    for key in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3",
                "accent4", "accent5", "accent6", "hlink", "folHlink"):
        node = etree.SubElement(scheme, f"{ns}{key}")
        etree.SubElement(node, f"{ns}srgbClr").set("val", _THEME_CLR[key])
    return scheme


def _font_scheme(head=Font.HEAD_THEME, body=Font.BODY_THEME):
    ns = f"{{{_A}}}"
    scheme = etree.Element(f"{ns}fontScheme")
    scheme.set("name", "Aqmen")
    for tag, face in (("majorFont", head), ("minorFont", body)):
        grp = etree.SubElement(scheme, f"{ns}{tag}")
        etree.SubElement(grp, f"{ns}latin").set("typeface", face)
        etree.SubElement(grp, f"{ns}ea").set("typeface", "")
        etree.SubElement(grp, f"{ns}cs").set("typeface", "")
    return scheme


def _apply_theme(prs):
    """Swap the presentation's colour + font schemes for the aqmen brand."""
    theme_part = prs.slide_masters[0].part.part_related_by(_RT.THEME)
    theme = etree.fromstring(theme_part.blob)
    elements = theme.find(_qn("a:themeElements"))
    elements.replace(elements.find(_qn("a:clrScheme")), _clr_scheme())
    elements.replace(elements.find(_qn("a:fontScheme")), _font_scheme())
    theme_part._blob = etree.tostring(theme, xml_declaration=True,
                                      encoding="UTF-8", standalone=True)


def _add_master_wordmark(prs, text="Aqmen"):
    """Bake the wordmark onto the slide master, bottom-left, brand navy."""
    master = prs.slide_masters[0]
    sp = master.shapes._spTree.add_textbox(100, "Wordmark", Inches(0.4),
                                           Inches(7.06), Inches(1.8),
                                           Inches(0.34))
    box = _Shape(sp, master.shapes)
    tf = box.text_frame
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = Font.HEAD
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = Palette.NAVY


def _clear_slides(prs):
    """Remove any slides from a presentation, keeping theme/master/layouts —
    so a populated template can be used as a clean base."""
    id_lst = prs.slides._sldIdLst
    for sldId in list(id_lst):
        rId = sldId.get(_qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        id_lst.remove(sldId)


# --------------------------------------------------------------------------- #
# Content value objects                                                        #
# --------------------------------------------------------------------------- #


@dataclass
class Bullet:
    """A single body bullet. `level` 0-2 controls indent; `bold` emphasises."""

    text: str
    level: int = 0
    bold: bool = False
    color: RGBColor | None = None


@dataclass
class ChartSpec:
    """A native chart. `kind` in {column, bar, line, stacked_column, stacked_bar}."""

    kind: str
    categories: Sequence[str]
    series: Sequence[tuple[str, Sequence[float]]]  # (name, values)
    number_format: str = "#,##0"
    data_labels: bool = True
    legend: bool = True
    y_title: str | None = None


@dataclass
class MekkoColumn:
    """One variable-width column of a Marimekko chart."""

    label: str
    width: float                       # relative width (e.g. a share or size)
    segments: Sequence[tuple[str, float]]  # (segment label, height value)


@dataclass
class HarveyRow:
    """One row of a harvey-ball comparison matrix. `fills` are 0.0–1.0 per column."""

    label: str
    cells: Sequence[float]


@dataclass
class DriverNode:
    """One box in a value driver / expression tree. `value` is the small strip
    under the label (e.g. "$mn", "[deals]"); `expr` is the formula shown muted
    under the label (e.g. "= deals × parties × penetration"); `certainty` 0/1/2 →
    low/medium/high dot. `parent` indexes the node in the previous column this one
    decomposes; `operator` (e.g. "×", "+") is how THIS node's children combine."""

    label: str
    value: str | None = None
    certainty: int | None = None
    expr: str | None = None
    parent: int = 0
    operator: str | None = None


@dataclass
class DriverColumn:
    """One level of a value driver / expression tree: a header + a stack of nodes."""

    header: str
    nodes: Sequence[DriverNode]


@dataclass
class MatrixItem:
    """An archetype region on a positioning matrix, drawn as the **bounding box**
    of its players' positions (+ padding), labelled — player names live in the
    revenue build, not in the box. `points` are member (x, y) positions in 0–1
    (left→right, bottom→top); the box spans their extent and small unlabelled dots
    mark each. `pad` grows the box beyond the points. If `points` is omitted, the
    box is centred at `x`/`y` sized by `size` (legacy single-box mode). `num` is a
    badge label."""

    label: str
    x: float = 0.5
    y: float = 0.5
    size: float = 0.5
    num: str | None = None
    points: Sequence[tuple[float, float]] | None = None
    pad: float = 0.05
    show_points: bool = True


@dataclass
class RangeColumn:
    """A column of a range grid. `kind`: "text" (plain), "num" (right-aligned
    number), "range" (a low–high bar with a point marker), or "balls" (harvey)."""

    name: str
    kind: str = "text"
    fmt: str = "{}"


@dataclass
class RangeRow:
    """A row of a range grid. `group` starts a new archetype band when it changes.
    `cells` align to the columns: str/number for text|num, (lo, pt, hi) for range,
    float 0–1 for balls."""

    label: str
    cells: Sequence
    group: str | None = None


_CHART_KINDS = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
}

# Valid data-label position per chart kind. OUTSIDE_END is ONLY legal on
# clustered column/bar — using it on a stacked chart makes PowerPoint report a
# content problem and offer to repair. Stacked labels go CENTER; line goes ABOVE.
_LABEL_POS = {
    "column": XL_LABEL_POSITION.OUTSIDE_END,
    "bar": XL_LABEL_POSITION.OUTSIDE_END,
    "line": XL_LABEL_POSITION.ABOVE,
    "stacked_column": XL_LABEL_POSITION.CENTER,
    "stacked_bar": XL_LABEL_POSITION.CENTER,
}


# --------------------------------------------------------------------------- #
# The Deck                                                                     #
# --------------------------------------------------------------------------- #


class Deck:
    """Builds an aqmen-house-style .pptx. One instance == one deck."""

    def __init__(self, draft: bool = True, wordmark: str = "Aqmen",
                 template: str | None = None):
        """Create a fully on-brand deck.

        With no `template`, the deck self-brands from blank — the aqmen theme
        (brand colours + Montserrat) is injected into the theme XML and the
        wordmark is baked onto the slide master, all in code. Pass a `template`
        (e.g. one of the bundled `*-template.pptx` starter decks, or your own
        branded base) to build on top of it instead; its existing slides are
        cleared so the deck starts clean while its theme, master, and layouts
        are kept."""
        if template and os.path.exists(template):
            self.prs = Presentation(template)
            self.prs.slide_width = SLIDE_W
            self.prs.slide_height = SLIDE_H
            _clear_slides(self.prs)
        else:
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_W
            self.prs.slide_height = SLIDE_H
            _apply_theme(self.prs)
            _add_master_wordmark(self.prs, wordmark)
        self._master_wordmark = True  # theme master carries the wordmark
        self._blank = self.prs.slide_layouts[6]
        self.draft = draft
        self.wordmark = wordmark
        self._page = 0  # incremented for every non-title slide

    def _suppress_master(self, slide):
        """Hide inherited master shapes (e.g. the wordmark) on a slide that
        draws its own chrome — the cover and full-bleed dividers."""
        slide._element.set("showMasterSp", "0")

    # ---- low-level helpers ------------------------------------------------- #

    def _slide(self):
        return self.prs.slides.add_slide(self._blank)

    def _text(
        self,
        slide,
        x,
        y,
        w,
        h,
        runs,
        *,
        size=12,
        color=Palette.INK,
        font=Font.BODY,
        bold=False,
        italic=False,
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
        line_spacing=1.05,
        space_after=2,
        wrap=True,
    ):
        """Add a textbox. `runs` is a string, or a list of paragraphs where each
        paragraph is a string or a list of (text, overrides-dict) run tuples."""
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)

        paragraphs = runs if isinstance(runs, list) else [runs]
        for i, para in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            p.space_after = Pt(space_after)
            p.space_before = Pt(0)
            run_specs = para if isinstance(para, list) else [(para, {})]
            for text, ov in run_specs:
                r = p.add_run()
                r.text = text
                r.font.size = Pt(ov.get("size", size))
                r.font.name = ov.get("font", font)
                r.font.bold = ov.get("bold", bold)
                r.font.italic = ov.get("italic", italic)
                r.font.color.rgb = ov.get("color", color)
        return box

    def _rect(self, slide, x, y, w, h, fill, line=None, line_w=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = line_w or Pt(0.75)
        shp.shadow.inherit = False
        return shp

    def _line(self, slide, x, y, w, color=Palette.GREY_LINE, weight=0.75):
        shp = slide.shapes.add_connector(2, x, y, x + w, y)  # 2 = straight
        shp.line.color.rgb = color
        shp.line.width = Pt(weight)
        return shp

    def _chrome(self, slide, eyebrow=None, source=None, page=True):
        """Standard slide furniture: DRAFT tag, eyebrow, wordmark, source, page #."""
        if self.draft:
            self._text(slide, Inches(0.2), Inches(0.1), Inches(2), Inches(0.3),
                       "DRAFT", size=9, color=Palette.AZURE, font=Font.HEAD,
                       bold=True)
        if eyebrow:
            head, sub = eyebrow if isinstance(eyebrow, tuple) else (eyebrow, None)
            runs = [(head + (":  " if sub else ""), {"bold": True})]
            if sub:
                runs.append((sub, {"italic": True}))
            self._text(slide, SLIDE_W - Inches(6.5) - Inches(0.2), Inches(0.12),
                       Inches(6.5), Inches(0.3), [runs], size=10.5,
                       color=Palette.NAVY, font=Font.HEAD, align=PP_ALIGN.RIGHT)
        # wordmark bottom-left — only if the template master doesn't provide it
        if not self._master_wordmark:
            self._text(slide, Inches(0.4), Inches(7.06), Inches(1.6),
                       Inches(0.32), self.wordmark, size=15, color=Palette.NAVY,
                       font=Font.HEAD, bold=True)
        if source:
            self._text(slide, Inches(1.55), Inches(7.14), Inches(9), Inches(0.28),
                       f"Source: {source}", size=8, color=Palette.GREY,
                       font=Font.BODY)
        if page:
            self._text(slide, SLIDE_W - Inches(0.9), Inches(7.12), Inches(0.6),
                       Inches(0.28), str(self._page), size=9, color=Palette.GREY,
                       align=PP_ALIGN.RIGHT)

    def _headline(self, slide, headline, illustrative=False):
        """The big navy headline band near the top of a content slide."""
        self._text(slide, MARGIN, Inches(0.62), SLIDE_W - 2 * MARGIN, Inches(1.15),
                   headline, size=24, color=Palette.NAVY, font=Font.HEAD,
                   bold=True, line_spacing=1.02, anchor=MSO_ANCHOR.TOP)
        if illustrative:
            self._tag(slide, "ILLUSTRATIVE", SLIDE_W - Inches(2.0), Inches(1.68))

    def _tag(self, slide, text, x, y):
        """A small boxed caption tag, e.g. ILLUSTRATIVE / NON-EXHAUSTIVE."""
        box = self._text(slide, x, y, Inches(1.7), Inches(0.24),
                         text, size=8.5, color=Palette.INK, font=Font.BODY,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        box.line.color.rgb = Palette.GREY
        box.line.width = Pt(0.5)
        return box

    def _rail(self, slide, takeaways: Sequence[str], title="Key Takeaways"):
        """The right-hand 'Key Takeaways' rail with a play-triangle header."""
        # vertical divider
        div = slide.shapes.add_connector(2, RAIL_X - Inches(0.3), CONTENT_TOP - Inches(0.35),
                                         RAIL_X - Inches(0.3), CONTENT_BOTTOM)
        div.line.color.rgb = Palette.GREY_LINE
        div.line.width = Pt(1)
        # play triangle
        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, RAIL_X,
                                     CONTENT_TOP - Inches(0.32), Inches(0.22),
                                     Inches(0.22))
        tri.rotation = 90
        tri.fill.solid()
        tri.fill.fore_color.rgb = Palette.NAVY
        tri.line.fill.background()
        tri.shadow.inherit = False
        self._text(slide, RAIL_X + Inches(0.32), CONTENT_TOP - Inches(0.34),
                   RAIL_W - Inches(0.32), Inches(0.3), title, size=13,
                   color=Palette.NAVY, font=Font.HEAD, bold=True)
        # bullets
        paras = [[(f"•  ", {"color": Palette.NAVY, "bold": True}),
                  (t, {})] for t in takeaways]
        self._text(slide, RAIL_X, CONTENT_TOP + Inches(0.15), RAIL_W,
                   CONTENT_BOTTOM - CONTENT_TOP - Inches(0.15), paras,
                   size=11, color=Palette.INK, line_spacing=1.08, space_after=7)

    def _body_bullets(self, slide, bullets: Sequence[Bullet], x, y, w, h):
        indents = {0: 0.0, 1: 0.28, 2: 0.56}
        marks = {0: "•", 1: "–", 2: "•"}
        paras = []
        for b in bullets:
            pad = "    " * b.level
            paras.append([
                (f"{pad}{marks.get(b.level, '•')}  ",
                 {"color": Palette.NAVY, "bold": True}),
                (b.text, {"bold": b.bold,
                          "color": b.color or Palette.INK}),
            ])
        self._text(slide, x, y, w, h, paras, size=12, color=Palette.INK,
                   line_spacing=1.1, space_after=6)

    # ---- public slide builders -------------------------------------------- #

    def title_slide(self, title, subtitle=None, date=None, kicker="Discussion Materials"):
        """Cover slide. Not paginated."""
        slide = self._slide()
        self._suppress_master(slide)
        self._rect(slide, Inches(0), Inches(0), Inches(0.28), SLIDE_H, Palette.AZURE)
        if self.draft:
            self._text(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.3),
                       "DRAFT", size=10, color=Palette.AZURE, font=Font.HEAD,
                       bold=True)
        self._text(slide, Inches(0.9), Inches(2.4), Inches(11), Inches(2.2),
                   title, size=38, color=Palette.NAVY, font=Font.HEAD, bold=True,
                   line_spacing=1.05)
        if subtitle:
            self._text(slide, Inches(0.9), Inches(4.5), Inches(10.5), Inches(0.8),
                       subtitle, size=20, color=Palette.BLUE, font=Font.HEAD)
        foot = kicker + (f"  |  {date}" if date else "")
        self._text(slide, Inches(0.9), Inches(6.6), Inches(8), Inches(0.4),
                   foot, size=13, color=Palette.INK, font=Font.BODY)
        self._text(slide, SLIDE_W - Inches(2.2), Inches(6.55), Inches(1.8),
                   Inches(0.4), self.wordmark, size=20, color=Palette.NAVY,
                   font=Font.HEAD, bold=True, align=PP_ALIGN.RIGHT)
        return slide

    def agenda(self, items: Sequence[str], active: str | None = None,
               subitems: dict | None = None, active_sub: str | None = None,
               title="Agenda"):
        """Agenda / section divider. `active` item is highlighted navy; an active
        item's `subitems` are listed below it, with `active_sub` shown bold."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, page=True)
        self._text(slide, MARGIN, Inches(0.62), Inches(8), Inches(0.9), title,
                   size=30, color=Palette.NAVY, font=Font.HEAD, bold=True)
        # centered list block
        lx, lw = Inches(4.4), Inches(4.6)
        y = Inches(2.2)
        self._line(slide, lx, y - Inches(0.12), lw, color=Palette.NAVY, weight=1)
        row_h = Inches(0.52)
        for it in items:
            is_active = (it == active)
            if is_active:
                self._rect(slide, lx, y, lw, row_h, Palette.NAVY_DEEP)
            self._text(slide, lx + Inches(0.15), y, lw - Inches(0.3), row_h, it,
                       size=15, font=Font.HEAD,
                       color=Palette.WHITE if is_active else Palette.NAVY,
                       bold=is_active, anchor=MSO_ANCHOR.MIDDLE)
            y = Emu(y + row_h)
            if is_active and subitems and it in subitems:
                for sub in subitems[it]:
                    is_as = (sub == active_sub)
                    if is_as:
                        self._rect(slide, lx, y, lw, row_h, Palette.CYAN_MUTED)
                    self._text(slide, lx + Inches(0.45), y, lw - Inches(0.6),
                               row_h, f"- {sub}", size=13, font=Font.HEAD,
                               color=Palette.NAVY, bold=is_as,
                               anchor=MSO_ANCHOR.MIDDLE)
                    y = Emu(y + row_h)
        self._line(slide, lx, y + Inches(0.08), lw, color=Palette.NAVY, weight=1)
        return slide

    def section_divider(self, section: str, subtitle: str | None = None):
        """A plain full-bleed section break slide."""
        self._page += 1
        slide = self._slide()
        self._suppress_master(slide)
        self._rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, Palette.NAVY_DEEP)
        self._text(slide, Inches(0.9), Inches(3.0), Inches(11), Inches(1.5),
                   section, size=34, color=Palette.WHITE, font=Font.HEAD, bold=True)
        if subtitle:
            self._text(slide, Inches(0.9), Inches(4.4), Inches(11), Inches(0.8),
                       subtitle, size=18, color=Palette.CYAN, font=Font.HEAD)
        self._text(slide, Inches(0.4), Inches(7.06), Inches(2), Inches(0.32),
                   self.wordmark, size=15, color=Palette.WHITE, font=Font.HEAD,
                   bold=True)
        return slide

    def executive_summary(self, rows: Sequence[tuple[str, Sequence[str]]],
                          headline="EXECUTIVE SUMMARY"):
        """Banded exec-summary matrix: left label column + right bullet blocks."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, page=True)
        self._text(slide, MARGIN, Inches(0.55), SLIDE_W - 2 * MARGIN, Inches(0.9),
                   headline, size=28, color=Palette.NAVY, font=Font.HEAD, bold=True)
        top = Inches(1.75)
        avail = CONTENT_BOTTOM - top
        n = len(rows)
        row_h = Emu(int(avail / n)) if n else avail
        label_w = Inches(2.3)
        body_x = MARGIN + label_w + Inches(0.2)
        body_w = SLIDE_W - MARGIN - body_x
        y = top
        for i, (label, bullets) in enumerate(rows):
            if i > 0:
                self._line(slide, MARGIN, y, SLIDE_W - 2 * MARGIN,
                           color=Palette.GREY_LINE, weight=0.75)
            self._text(slide, MARGIN, y + Inches(0.1), label_w, row_h, label,
                       size=14, color=Palette.NAVY, font=Font.HEAD, bold=True,
                       anchor=MSO_ANCHOR.TOP)
            paras = [[(f"•  ", {"color": Palette.NAVY, "bold": True}),
                      (b, {})] for b in bullets]
            self._text(slide, body_x, y + Inches(0.08), body_w, row_h, paras,
                       size=10.5, color=Palette.INK, line_spacing=1.05,
                       space_after=3)
            y = Emu(y + row_h)
        return slide

    def content_slide(self, headline, body: Sequence[Bullet] | None = None,
                      takeaways: Sequence[str] | None = None,
                      eyebrow=None, source=None, illustrative=False,
                      left_title: str | None = None):
        """The workhorse: headline + left body bullets + optional right takeaways
        rail. If `takeaways` is None the body spans the full content width."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)
        body_w = (RAIL_X - Inches(0.55) - MARGIN) if takeaways else (SLIDE_W - 2 * MARGIN)
        y = CONTENT_TOP
        if left_title:
            self._text(slide, MARGIN, y, body_w, Inches(0.35), left_title,
                       size=14, color=Palette.NAVY, font=Font.HEAD, bold=True)
            y = Emu(y + Inches(0.45))
        if body:
            self._body_bullets(slide, body, MARGIN, y, body_w,
                               CONTENT_BOTTOM - y)
        if takeaways:
            self._rail(slide, takeaways)
        return slide

    def chart_slide(self, headline, chart: ChartSpec,
                    takeaways: Sequence[str] | None = None,
                    eyebrow=None, source=None, illustrative=False,
                    chart_title: str | None = None):
        """Headline + a native chart (left) + optional takeaways rail (right)."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)
        cx, cy = MARGIN, CONTENT_TOP + (Inches(0.35) if chart_title else Inches(0))
        cw = (RAIL_X - Inches(0.55) - MARGIN) if takeaways else (SLIDE_W - 2 * MARGIN)
        ch = CONTENT_BOTTOM - cy
        if chart_title:
            self._text(slide, MARGIN, CONTENT_TOP - Inches(0.1), cw, Inches(0.35),
                       chart_title, size=14, color=Palette.NAVY, font=Font.HEAD,
                       bold=True)
        self._add_chart(slide, chart, cx, cy, cw, ch)
        if takeaways:
            self._rail(slide, takeaways)
        return slide

    def _add_chart(self, slide, spec: ChartSpec, x, y, w, h):
        ct = _CHART_KINDS.get(spec.kind)
        if ct is None:
            raise ValueError(f"unknown chart kind: {spec.kind!r}")
        data = CategoryChartData()
        data.categories = list(spec.categories)
        for name, vals in spec.series:
            data.add_series(name, tuple(vals), spec.number_format)
        gframe = slide.shapes.add_chart(ct, x, y, w, h, data)
        chart = gframe.chart
        chart.has_title = False
        # legend
        if spec.legend and len(spec.series) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = Font.BODY
        else:
            chart.has_legend = False
        # colour series from the brand palette
        for i, plot_series in enumerate(chart.series):
            color = Palette.SERIES[i % len(Palette.SERIES)]
            fmt = plot_series.format
            if spec.kind == "line":
                fmt.line.color.rgb = color
                fmt.line.width = Pt(2.25)
            else:
                fmt.fill.solid()
                fmt.fill.fore_color.rgb = color
                fmt.line.fill.background()
            if spec.data_labels:
                plot_series.data_labels.number_format = spec.number_format
                plot_series.data_labels.number_format_is_linked = False
                plot_series.data_labels.font.size = Pt(9)
                plot_series.data_labels.font.name = Font.BODY
                # Label position must be valid for the chart type, or PowerPoint
                # flags the file and offers to "repair" it. OUTSIDE_END is only
                # allowed on clustered column/bar; stacked must use CENTER; line
                # uses ABOVE.
                pos = _LABEL_POS.get(spec.kind)
                if pos is not None:
                    try:
                        plot_series.data_labels.position = pos
                    except Exception:
                        pass
        # axes styling
        try:
            for axis in (chart.category_axis, chart.value_axis):
                axis.tick_labels.font.size = Pt(9)
                axis.tick_labels.font.name = Font.BODY
                axis.format.line.color.rgb = Palette.GREY_LINE
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = Palette.GREY_BG
            if spec.y_title:
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = spec.y_title
        except Exception:
            pass
        return chart

    def marimekko_slide(self, headline, columns: Sequence[MekkoColumn],
                        takeaways: Sequence[str] | None = None,
                        eyebrow=None, source=None, illustrative=True,
                        chart_title: str | None = None, value_fmt="{:.1f}"):
        """A variable-width stacked ('Mekko'/Marimekko) chart drawn from shapes —
        column widths encode one dimension, stack heights another. This is the
        flagship market-sizing visual (e.g. TAM by region x SAM/whitespace)."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)
        if chart_title:
            self._text(slide, MARGIN, CONTENT_TOP - Inches(0.15),
                       SLIDE_W - 2 * MARGIN, Inches(0.35), chart_title, size=14,
                       color=Palette.NAVY, font=Font.HEAD, bold=True)

        area_x = MARGIN
        area_y = CONTENT_TOP + Inches(1.05)
        area_w = (RAIL_X - Inches(0.55) - MARGIN) if takeaways else Inches(11.0)
        area_h = CONTENT_BOTTOM - area_y - Inches(0.35)

        seg_labels = [s[0] for s in columns[0].segments]
        seg_colors, ci = {}, 0
        for lbl in seg_labels:
            if "whitespace" in lbl.lower():
                seg_colors[lbl] = Palette.CYAN_MUTED   # hatched at draw time
            else:
                seg_colors[lbl] = Palette.SERIES[ci % len(Palette.SERIES)]
                ci += 1
        # legend on its own row between the chart title and the columns
        self._legend(slide, [(lbl, seg_colors[lbl]) for lbl in seg_labels],
                     area_x, CONTENT_TOP + Inches(0.32))
        total_w = sum(c.width for c in columns) or 1.0
        max_h = max(sum(v for _, v in c.segments) for c in columns) or 1.0
        gap = Inches(0.12)
        usable_w = Emu(int(area_w - gap * (len(columns) - 1)))

        cursor_x = area_x
        for col in columns:
            col_w = Emu(int(usable_w * (col.width / total_w)))
            col_total = sum(v for _, v in col.segments)
            # column header (label + share)
            self._text(slide, cursor_x, area_y - Inches(0.5), col_w, Inches(0.45),
                       [[(f"{value_fmt.format(col_total)}", {"bold": True,
                          "color": Palette.NAVY, "size": 12})],
                        [(f"({col.width / total_w * 100:.0f}%)",
                          {"italic": True, "color": Palette.GREY, "size": 10})]],
                       align=PP_ALIGN.CENTER)
            seg_y = area_y
            for lbl, val in col.segments:
                seg_h = Emu(int(area_h * (val / max_h)))
                is_ws = "whitespace" in lbl.lower()
                shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cursor_x, seg_y,
                                             col_w, seg_h)
                shp.shadow.inherit = False
                shp.line.color.rgb = Palette.WHITE
                shp.line.width = Pt(1)
                if is_ws:
                    try:
                        shp.fill.patterned()
                        shp.fill.pattern = MSO_PATTERN_TYPE.LIGHT_DOWNWARD_DIAGONAL
                        shp.fill.fore_color.rgb = Palette.STEEL
                        shp.fill.back_color.rgb = Palette.WHITE
                    except Exception:
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = Palette.CYAN_MUTED
                    txt_color = Palette.NAVY
                else:
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = seg_colors[lbl]
                    txt_color = Palette.WHITE
                self._text(slide, cursor_x, seg_y, col_w, seg_h,
                           value_fmt.format(val), size=10, color=txt_color,
                           font=Font.BODY, bold=True, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE)
                seg_y = Emu(seg_y + seg_h)
            # x label under the column
            self._text(slide, cursor_x, area_y + area_h + Inches(0.05), col_w,
                       Inches(0.3), col.label, size=10.5, color=Palette.NAVY,
                       font=Font.HEAD, bold=True, align=PP_ALIGN.CENTER)
            cursor_x = Emu(cursor_x + col_w + gap)

        if takeaways:
            self._rail(slide, takeaways)
        return slide

    def _legend(self, slide, entries: Sequence[tuple[str, RGBColor]], x, y):
        cur = x
        for lbl, color in entries:
            self._rect(slide, cur, y, Inches(0.16), Inches(0.16), color)
            self._text(slide, cur + Inches(0.24), y - Inches(0.03), Inches(2.2),
                       Inches(0.25), lbl, size=9.5, color=Palette.INK,
                       font=Font.BODY, anchor=MSO_ANCHOR.MIDDLE)
            # advance: swatch + gap + estimated text width (~0.075in per char)
            cur = Emu(cur + Inches(0.24) + Inches(0.075 * len(lbl) + 0.15)
                      + Inches(0.35))

    def harvey_matrix_slide(self, headline, columns: Sequence[str],
                            rows: Sequence[HarveyRow], row_desc: dict | None = None,
                            takeaways: Sequence[str] | None = None,
                            eyebrow=None, source=None, illustrative=False,
                            row_label_header="", left_w=2.4):
        """A comparison matrix with harvey balls (0=empty … 1=full) — the standard
        competitive-benchmarking visual. `columns` are compared entities; each
        `HarveyRow` gives a criterion label and one fill (0–1) per column."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)

        table_w = (RAIL_X - Inches(0.55) - MARGIN) if takeaways else (SLIDE_W - 2 * MARGIN)
        left = Inches(left_w)
        desc_w = Inches(2.6) if row_desc else Inches(0)
        top = CONTENT_TOP + Inches(0.1)
        n_cols = len(columns)
        col_area = Emu(int(table_w - left - desc_w))
        col_w = Emu(int(col_area / n_cols)) if n_cols else col_area
        n_rows = len(rows)
        row_h = Emu(int((CONTENT_BOTTOM - top - Inches(0.5)) / max(n_rows, 1)))

        # header row
        self._text(slide, MARGIN, top, left, Inches(0.5), row_label_header,
                   size=11, color=Palette.NAVY, font=Font.HEAD, bold=True,
                   anchor=MSO_ANCHOR.BOTTOM)
        for j, col in enumerate(columns):
            cx = Emu(MARGIN + left + desc_w + col_w * j)
            self._text(slide, cx, top, col_w, Inches(0.5), col, size=10.5,
                       color=Palette.NAVY, font=Font.HEAD, bold=True,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        y = Emu(top + Inches(0.55))
        for r in rows:
            self._line(slide, MARGIN, y, table_w, color=Palette.GREY_LINE, weight=0.5)
            self._text(slide, MARGIN, y, left, row_h, r.label, size=11,
                       color=Palette.NAVY, font=Font.HEAD, bold=True,
                       anchor=MSO_ANCHOR.MIDDLE)
            if row_desc and r.label in row_desc:
                self._text(slide, MARGIN + left, y, desc_w, row_h,
                           row_desc[r.label], size=9.5, color=Palette.INK,
                           anchor=MSO_ANCHOR.MIDDLE)
            for j, fill in enumerate(r.cells):
                cx = Emu(MARGIN + left + desc_w + col_w * j)
                self._harvey_ball(slide, Emu(cx + col_w / 2 - Inches(0.12)),
                                  Emu(y + row_h / 2 - Inches(0.12)),
                                  Inches(0.24), fill)
            y = Emu(y + row_h)
        self._line(slide, MARGIN, y, table_w, color=Palette.GREY_LINE, weight=0.5)
        if takeaways:
            self._rail(slide, takeaways)
        return slide

    def _harvey_ball(self, slide, x, y, d, fill: float):
        """A harvey ball: navy ring with a `fill` (0–1) navy pie wedge."""
        ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        ring.fill.solid()
        ring.fill.fore_color.rgb = Palette.WHITE
        ring.line.color.rgb = Palette.NAVY
        ring.line.width = Pt(1.25)
        ring.shadow.inherit = False
        fill = max(0.0, min(1.0, fill))
        if fill <= 0:
            return
        if fill >= 1:
            wedge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        else:
            wedge = slide.shapes.add_shape(MSO_SHAPE.PIE, x, y, d, d)
            # PIE default sweeps; set adjustments: start at top (270°), sweep cw
            try:
                wedge.adjustments[0] = 270
                wedge.adjustments[1] = (270 + 360 * fill) % 360
            except Exception:
                pass
        wedge.fill.solid()
        wedge.fill.fore_color.rgb = Palette.NAVY
        wedge.line.fill.background()
        wedge.shadow.inherit = False

    # ---- value driver tree ------------------------------------------------- #

    def _rounded(self, slide, x, y, w, h, fill, line=None, line_w=None, radius=0.08):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = line_w or Pt(0.75)
        shp.shadow.inherit = False
        return shp

    def _driver_node(self, slide, x, y, w, h, node: DriverNode):
        """A two-tone driver box: label (+ optional muted expression) above a navy
        value strip, with an optional certainty dot."""
        self._rounded(slide, x, y, w, h, Palette.CYAN_MUTED,
                      line=Palette.AZURE, line_w=Pt(0.75))
        iw = Emu(w - Inches(0.12))
        ix = Emu(x + Inches(0.06))
        strip_h = Emu(min(int(h * 0.36), int(Inches(0.30)))) if node.value is not None else Emu(0)
        expr_h = Inches(0.2) if node.expr else Emu(0)
        label_h = Emu(h - strip_h - expr_h)
        self._text(slide, ix, y, iw, label_h, node.label, size=9.5,
                   color=Palette.NAVY, font=Font.HEAD, bold=True,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        if node.expr:
            self._text(slide, ix, Emu(y + label_h), iw, expr_h, node.expr,
                       size=7.5, color=Palette.BLUE, font=Font.BODY, italic=True,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if node.value is not None:
            self._rounded(slide, ix, Emu(y + label_h + expr_h - Inches(0.02)),
                          iw, strip_h, Palette.NAVY_DEEP, radius=0.12)
            self._text(slide, ix, Emu(y + label_h + expr_h - Inches(0.02)), iw,
                       strip_h, node.value, size=9, color=Palette.WHITE,
                       font=Font.BODY, bold=True, align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
        if node.certainty is not None:
            c = Palette.CERTAINTY[max(0, min(2, node.certainty))]
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Emu(x + w - Inches(0.17)),
                                         Emu(y + Inches(0.05)), Inches(0.12),
                                         Inches(0.12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = c
            dot.line.color.rgb = Palette.WHITE
            dot.line.width = Pt(0.75)
            dot.shadow.inherit = False

    def _operator_glyph(self, slide, cx, cy, text):
        d = Inches(0.3)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx - d // 2),
                                      Emu(cy - d // 2), d, d)
        circ.fill.solid()
        circ.fill.fore_color.rgb = Palette.WHITE
        circ.line.color.rgb = Palette.NAVY
        circ.line.width = Pt(1)
        circ.shadow.inherit = False
        self._text(slide, Emu(cx - d // 2), Emu(cy - d // 2), d, d, text,
                   size=13, color=Palette.NAVY, font=Font.HEAD, bold=True,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    def driver_tree_slide(self, headline, columns: Sequence[DriverColumn],
                          takeaways: Sequence[str] | None = None,
                          eyebrow=None, source=None, illustrative=True,
                          chart_title: str | None = None, certainty_key=True,
                          note: str | None = None):
        """A left→right **expression tree**: the market expression decomposed into
        its variables and their drivers, boxes joined to their parent by a grouped
        connector spine, with per-node ×/+ operators and certainty dots. Keep it to
        the variables/expressions — note segmentation compactly via `note` rather
        than exploding every segment box. The flagship market-sizing visual."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)
        if chart_title:
            self._text(slide, MARGIN, CONTENT_TOP - Inches(0.15),
                       SLIDE_W - 2 * MARGIN, Inches(0.35), chart_title, size=14,
                       color=Palette.NAVY, font=Font.HEAD, bold=True)
        if certainty_key:
            # compact horizontal legend on the chart-title row, right of centre
            ky = CONTENT_TOP - Inches(0.12)
            kx = SLIDE_W - Inches(5.6)
            self._text(slide, Emu(kx - Inches(1.4)), ky, Inches(1.35), Inches(0.25),
                       "Certainty:", size=9.5, color=Palette.NAVY, font=Font.HEAD,
                       bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            for lbl, ci in (("Low", 0), ("Medium", 1), ("High", 2)):
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, kx,
                                             Emu(ky + Inches(0.04)),
                                             Inches(0.13), Inches(0.13))
                dot.fill.solid()
                dot.fill.fore_color.rgb = Palette.CERTAINTY[ci]
                dot.line.fill.background()
                dot.shadow.inherit = False
                self._text(slide, Emu(kx + Inches(0.18)), ky, Inches(0.85),
                           Inches(0.25), lbl, size=9.5, color=Palette.INK,
                           font=Font.BODY, anchor=MSO_ANCHOR.MIDDLE)
                kx = Emu(kx + Inches(1.15))
            takeaways = None  # tree spans full width; no rail

        area_x = MARGIN
        area_y = CONTENT_TOP + Inches(0.55)
        area_w = SLIDE_W - 2 * MARGIN
        area_h = CONTENT_BOTTOM - area_y - Inches(0.15)
        n = len(columns)
        gap = Inches(0.5)                       # room for the operator spine
        col_w = Emu(int((area_w - gap * (n - 1)) / n))
        node_gap = Inches(0.18)

        centers = []  # (x_left, x_right, [node_center_y]) per column
        cursor_x = area_x
        for col in columns:
            # header
            self._text(slide, cursor_x, area_y - Inches(0.42), col_w, Inches(0.35),
                       col.header, size=10.5, color=Palette.NAVY, font=Font.HEAD,
                       bold=True, italic=True, align=PP_ALIGN.CENTER)
            m = len(col.nodes)
            nh = Emu(int((area_h - node_gap * (m - 1)) / m))
            nh = Emu(min(nh, int(Inches(0.95))))
            block_h = Emu(nh * m + node_gap * (m - 1))
            ny = Emu(area_y + (area_h - block_h) // 2)
            node_centers = []
            for node in col.nodes:
                self._driver_node(slide, cursor_x, ny, col_w, nh, node)
                node_centers.append(Emu(ny + nh // 2))
                ny = Emu(ny + nh + node_gap)
            centers.append((cursor_x, Emu(cursor_x + col_w), node_centers))
            cursor_x = Emu(cursor_x + col_w + gap)

        # connectors: each child links to its PARENT node in the previous column,
        # grouped per parent so sub-trees stay distinct (no misleading full bus).
        for i in range(1, n):
            lx = centers[i - 1][1]
            rx = centers[i][0]
            parent_ys = centers[i - 1][2]
            child_ys = centers[i][2]
            groups = {}
            for j, node in enumerate(columns[i].nodes):
                p = max(0, min(len(parent_ys) - 1, node.parent))
                groups.setdefault(p, []).append(child_ys[j])
            for p, ys in groups.items():
                p_cy = parent_ys[p]
                bus_x = Emu((lx + rx) // 2)
                top_y = min(ys + [p_cy])
                bot_y = max(ys + [p_cy])
                v = slide.shapes.add_connector(2, bus_x, top_y, bus_x, bot_y)
                v.line.color.rgb = Palette.GREY_LINE
                v.line.width = Pt(1)
                self._line(slide, lx, p_cy, Emu(bus_x - lx), color=Palette.GREY_LINE, weight=1)
                for yy in ys:
                    self._line(slide, bus_x, yy, Emu(rx - bus_x), color=Palette.GREY_LINE, weight=1)
                op = columns[i - 1].nodes[p].operator
                if op:
                    self._operator_glyph(slide, Emu((lx + bus_x) // 2), p_cy, op)

        if note:
            self._text(slide, MARGIN, Emu(CONTENT_BOTTOM - Inches(0.02)),
                       SLIDE_W - 2 * MARGIN, Inches(0.3),
                       [[("▸ ", {"color": Palette.NAVY, "bold": True}),
                         (note, {"italic": True})]], size=9,
                       color=Palette.GREY, font=Font.BODY)
        if takeaways:
            self._rail(slide, takeaways)
        return slide

    # ---- positioning matrix ------------------------------------------------ #

    def _num_badge(self, slide, x, y, d, text, fill=None):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        c.fill.solid()
        c.fill.fore_color.rgb = fill or Palette.WHITE
        c.line.color.rgb = Palette.NAVY
        c.line.width = Pt(1)
        c.shadow.inherit = False
        self._text(slide, x, y, d, d, text, size=9,
                   color=(Palette.WHITE if fill else Palette.NAVY),
                   font=Font.HEAD, bold=True, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE)

    def positioning_matrix_slide(self, headline, items: Sequence[MatrixItem],
                                 x_title, y_title, x_ticks=None, y_ticks=None,
                                 takeaways: Sequence[str] | None = None,
                                 eyebrow=None, source=None, illustrative=True,
                                 chart_title: str | None = None):
        """A 2-axis positioning matrix: labelled, size-scaled boxes placed on a
        plane (e.g. company size × service focus). The flagship competitive-
        landscape framing visual. `items` carry x/y in 0–1 and a size in 0–1."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)
        if chart_title:
            self._text(slide, MARGIN, CONTENT_TOP - Inches(0.15),
                       SLIDE_W - 2 * MARGIN, Inches(0.35), chart_title, size=14,
                       color=Palette.NAVY, font=Font.HEAD, bold=True)

        right = (RAIL_X - Inches(0.55)) if takeaways else (SLIDE_W - MARGIN)
        x0 = MARGIN + Inches(1.55)          # room for y-axis title + ticks
        x1 = right - Inches(0.2)
        y_top = CONTENT_TOP + Inches(0.35)
        y_bot = CONTENT_BOTTOM - Inches(0.55)
        pw, ph = Emu(x1 - x0), Emu(y_bot - y_top)

        # axes (arrows)
        xa = slide.shapes.add_connector(2, x0, y_bot, x1, y_bot)
        ya = slide.shapes.add_connector(2, x0, y_bot, x0, y_top)
        for a in (xa, ya):
            a.line.color.rgb = Palette.NAVY
            a.line.width = Pt(1.5)
            try:
                a.line.headEnd if False else None
            except Exception:
                pass
        # axis titles
        self._text(slide, x0, Emu(y_bot + Inches(0.28)), pw, Inches(0.3), x_title,
                   size=11, color=Palette.NAVY, font=Font.HEAD, bold=True,
                   align=PP_ALIGN.CENTER)
        # rotated 270°: a textbox pivots on its CENTRE, so place the centre where
        # the vertical label should read (just left of the y-axis, mid-plot).
        yw, yh = Inches(3.0), Inches(0.3)
        yc_x, yc_y = Emu(x0 - Inches(1.3)), Emu(y_top + ph // 2)
        yt = self._text(slide, Emu(yc_x - yw // 2), Emu(yc_y - yh // 2), yw, yh,
                        y_title, size=11, color=Palette.NAVY, font=Font.HEAD,
                        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        yt.rotation = 270
        # ticks
        for i, tk in enumerate(x_ticks or []):
            n = len(x_ticks)
            tx = Emu(x0 + pw * (i + 0.5) // n) if n else x0
            self._text(slide, Emu(tx - Inches(0.9)), Emu(y_bot + Inches(0.05)),
                       Inches(1.8), Inches(0.24), tk, size=9.5, color=Palette.INK,
                       align=PP_ALIGN.CENTER)
        for i, tk in enumerate(y_ticks or []):
            n = len(y_ticks)
            ty = Emu(y_bot - ph * (i + 0.5) // n) if n else y_bot
            self._text(slide, Emu(x0 - Inches(1.05)), Emu(ty - Inches(0.12)),
                       Inches(0.95), Inches(0.24), tk, size=9.5, color=Palette.INK,
                       align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

        def px(fx):
            return Emu(x0 + int(pw * max(0.0, min(1.0, fx))))

        def py(fy):
            return Emu(y_bot - int(ph * max(0.0, min(1.0, fy))))

        # archetype regions — each drawn as the bounding box of its players'
        # points (+ padding), labelled; small unlabelled dots mark the players.
        for it in items:
            pts = list(it.points or [])
            if pts:
                minx = max(0.0, min(p[0] for p in pts) - it.pad)
                maxx = min(1.0, max(p[0] for p in pts) + it.pad)
                miny = max(0.0, min(p[1] for p in pts) - it.pad)
                maxy = min(1.0, max(p[1] for p in pts) + it.pad)
            else:
                hw = 0.5 * it.size * 0.5
                minx, maxx = it.x - hw, it.x + hw
                miny, maxy = it.y - hw * 1.2, it.y + hw * 1.2
            bx, by = px(minx), py(maxy)
            bw = Emu(px(maxx) - px(minx))
            bh = Emu(py(miny) - py(maxy))
            self._rounded(slide, bx, by, bw, bh, Palette.CYAN_MUTED,
                          line=Palette.AZURE, line_w=Pt(1), radius=0.04)
            self._text(slide, Emu(bx + Inches(0.05)), Emu(by + Inches(0.04)),
                       Emu(bw - Inches(0.1)), Inches(0.24), it.label, size=9.5,
                       color=Palette.NAVY, font=Font.HEAD, bold=True,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            if it.show_points:
                for (fx, fy) in pts:
                    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                                 Emu(px(fx) - Inches(0.05)),
                                                 Emu(py(fy) - Inches(0.05)),
                                                 Inches(0.1), Inches(0.1))
                    dot.fill.solid()
                    dot.fill.fore_color.rgb = Palette.NAVY
                    dot.line.color.rgb = Palette.WHITE
                    dot.line.width = Pt(0.75)
                    dot.shadow.inherit = False
            if it.num:
                self._num_badge(slide, Emu(bx - Inches(0.1)), Emu(by - Inches(0.1)),
                                Inches(0.26), it.num, fill=Palette.NAVY)

        if takeaways:
            self._rail(slide, takeaways)
        return slide

    # ---- sensitivity heatmap ----------------------------------------------- #

    @staticmethod
    def _lerp(a: RGBColor, b: RGBColor, t: float) -> RGBColor:
        t = max(0.0, min(1.0, t))
        return RGBColor(int(a[0] + (b[0] - a[0]) * t),
                        int(a[1] + (b[1] - a[1]) * t),
                        int(a[2] + (b[2] - a[2]) * t))

    def heatmap_slide(self, headline, cols: Sequence[str], rows: Sequence[str],
                      values: Sequence[Sequence[float]], base=None,
                      col_header="", row_header="", value_fmt="{:.0f}",
                      takeaways: Sequence[str] | None = None, eyebrow=None,
                      source=None, illustrative=False, chart_title: str | None = None):
        """A sensitivity matrix: a `rows × cols` grid of cells shaded on a light→
        navy ramp by value, with the base-case cell outlined. The standard DCF
        WACC × terminal-growth (or any two-driver) sensitivity view. `base` is a
        (row, col) index tuple to highlight."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)
        if chart_title:
            self._text(slide, MARGIN, CONTENT_TOP - Inches(0.15),
                       SLIDE_W - 2 * MARGIN, Inches(0.35), chart_title, size=14,
                       color=Palette.NAVY, font=Font.HEAD, bold=True)

        right = (RAIL_X - Inches(0.55)) if takeaways else (SLIDE_W - MARGIN)
        left_w = Inches(1.35)
        top = CONTENT_TOP + Inches(0.75)
        grid_x = Emu(MARGIN + left_w)
        grid_w = Emu(right - grid_x)
        grid_h = Emu(CONTENT_BOTTOM - top - Inches(0.2))
        nr, nc = len(rows), len(cols)
        cw = Emu(grid_w // nc)
        ch = Emu(grid_h // nr)
        flat = [v for row in values for v in row]
        lo, hi = min(flat), max(flat)
        span = (hi - lo) or 1.0

        # axis titles
        self._text(slide, grid_x, Emu(top - Inches(0.62)), grid_w, Inches(0.26),
                   col_header, size=10.5, color=Palette.NAVY, font=Font.HEAD,
                   bold=True, align=PP_ALIGN.CENTER)
        self._text(slide, MARGIN, Emu(top - Inches(0.62)), left_w, Inches(0.26),
                   row_header, size=10.5, color=Palette.NAVY, font=Font.HEAD,
                   bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        # column headers
        for j, c in enumerate(cols):
            self._text(slide, Emu(grid_x + cw * j), Emu(top - Inches(0.3)), cw,
                       Inches(0.28), c, size=10, color=Palette.INK, font=Font.BODY,
                       bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # rows
        for i, r in enumerate(rows):
            ry = Emu(top + ch * i)
            self._text(slide, MARGIN, ry, left_w, ch, r, size=10,
                       color=Palette.INK, font=Font.BODY, bold=True,
                       align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            for j in range(nc):
                v = values[i][j]
                t = (v - lo) / span
                fill = self._lerp(Palette.CYAN, Palette.NAVY, t)
                is_base = base is not None and tuple(base) == (i, j)
                cell = self._rect(slide, Emu(grid_x + cw * j), ry, cw, ch, fill,
                                  line=(Palette.ORANGE if is_base else Palette.WHITE),
                                  line_w=Pt(2.5 if is_base else 1))
                self._text(slide, Emu(grid_x + cw * j), ry, cw, ch,
                           value_fmt.format(v), size=10.5,
                           color=(Palette.WHITE if t > 0.55 else Palette.INK),
                           font=Font.BODY, bold=is_base, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE)
        if takeaways:
            self._rail(slide, takeaways)
        return slide

    # ---- competitive-landscape revenue build ------------------------------ #

    def revenue_build_slide(self, headline, groups,
                            info_cols=(("hq", "HQ"), ("employees", "Employees")),
                            rev_key="revenue", pct_key="pct", mrev_key="market_rev",
                            unit="$m", total_label="Total addressable (bottom-up market)",
                            total=None, eyebrow=None, source=None, illustrative=False,
                            chart_title: str | None = None, value_fmt="{:,.0f}"):
        """The competitive-landscape flagship: a player **revenue build** — each
        player's total revenue × its % addressable to this market = the revenue
        assigned to the market, bar-visualised and summed bottom-up to a market
        total. `groups` is a list of (archetype_label, [player_dict, …]); each
        player_dict has `name`, the `info_cols` keys, `rev_key`, `pct_key`,
        `mrev_key`. Players with only revenue/market_rev render as revenue-only."""
        self._page += 1
        slide = self._slide()
        self._chrome(slide, eyebrow=eyebrow, source=source)
        self._headline(slide, headline, illustrative=illustrative)
        if chart_title:
            self._text(slide, MARGIN, CONTENT_TOP - Inches(0.15),
                       SLIDE_W - 2 * MARGIN, Inches(0.35), chart_title, size=14,
                       color=Palette.NAVY, font=Font.HEAD, bold=True)

        tw = SLIDE_W - 2 * MARGIN
        # column x-fractions: name | info… | revenue | % | (market-rev bar region)
        name_f = 0.22
        info_f = 0.11
        rev_f, pct_f = 0.12, 0.10
        bar_f = 1.0 - (name_f + info_f * len(info_cols) + rev_f + pct_f)
        xs = [0.0, name_f]
        for _ in info_cols:
            xs.append(xs[-1] + info_f)
        xs.append(xs[-1] + rev_f)      # after revenue
        xs.append(xs[-1] + pct_f)      # after pct → start of bar region
        colx = [Emu(MARGIN + int(tw * f)) for f in xs]
        bar_x0 = Emu(colx[-1] + Inches(0.15))            # gutter before bar region
        bar_track = Emu(int(tw * bar_f) - Inches(1.0))   # room for gutter + value label

        top = CONTENT_TOP + Inches(0.75)
        # header row
        heads = ["Company"] + [h for _, h in info_cols] + \
                [f"Revenue ({unit})", "% to market"]
        for i, h in enumerate(heads):
            al = PP_ALIGN.RIGHT if i >= 1 + len(info_cols) else PP_ALIGN.LEFT
            self._text(slide, colx[i], Emu(top - Inches(0.34)),
                       Emu(colx[i + 1] - colx[i]), Inches(0.3), h, size=9.5,
                       color=Palette.NAVY, font=Font.HEAD, bold=True, align=al,
                       anchor=MSO_ANCHOR.BOTTOM)
        self._text(slide, bar_x0, Emu(top - Inches(0.34)), Emu(tw - (bar_x0 - MARGIN)),
                   Inches(0.3), f"Market revenue ({unit})", size=9.5,
                   color=Palette.NAVY, font=Font.HEAD, bold=True,
                   align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)
        self._line(slide, MARGIN, top, tw, color=Palette.NAVY, weight=1)

        n_rows = sum(1 + len(pl) for _, pl in groups) + 1  # +total
        row_h = Emu(min(int((CONTENT_BOTTOM - top - Inches(0.1)) / max(n_rows, 1)),
                        int(Inches(0.34))))
        all_mrev = [p.get(mrev_key, 0) or 0 for _, pl in groups for p in pl]
        mx = max(all_mrev) or 1.0  # scale player bars to the largest player

        y = Emu(top + Inches(0.06))
        for arch, players in groups:
            self._rect(slide, MARGIN, y, tw, row_h, Palette.GREY_BG)
            self._text(slide, Emu(MARGIN + Inches(0.06)), y, tw, row_h, arch,
                       size=9.5, color=Palette.NAVY, font=Font.HEAD, bold=True,
                       anchor=MSO_ANCHOR.MIDDLE)
            y = Emu(y + row_h)
            for p in players:
                self._text(slide, Emu(colx[0] + Inches(0.06)), y,
                           Emu(colx[1] - colx[0]), row_h, p.get("name", ""),
                           size=9.5, color=Palette.INK, font=Font.BODY, bold=True,
                           anchor=MSO_ANCHOR.MIDDLE)
                for k, (key, _) in enumerate(info_cols):
                    self._text(slide, colx[1 + k], y, Emu(colx[2 + k] - colx[1 + k]),
                               row_h, str(p.get(key, "")), size=9, color=Palette.INK,
                               font=Font.BODY, anchor=MSO_ANCHOR.MIDDLE)
                ci = 1 + len(info_cols)
                rev = p.get(rev_key)
                if rev is not None:
                    self._text(slide, Emu(colx[ci] - Inches(0.12)), y,
                               Emu(colx[ci + 1] - colx[ci]), row_h,
                               value_fmt.format(rev), size=9, color=Palette.INK,
                               font=Font.BODY, align=PP_ALIGN.RIGHT,
                               anchor=MSO_ANCHOR.MIDDLE)
                pct = p.get(pct_key)
                if pct is not None:
                    self._text(slide, Emu(colx[ci + 1] - Inches(0.12)), y,
                               Emu(bar_x0 - colx[ci + 1]), row_h, f"{pct:g}%",
                               size=9, color=Palette.INK, font=Font.BODY,
                               align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
                mrev = p.get(mrev_key)
                if mrev is not None:
                    bw = Emu(max(int(bar_track * (mrev / mx)), int(Inches(0.02))))
                    self._rect(slide, bar_x0, Emu(y + row_h // 2 - Inches(0.07)),
                               bw, Inches(0.14), Palette.AZURE)
                    self._text(slide, Emu(bar_x0 + bw + Inches(0.05)), y,
                               Inches(0.9), row_h, value_fmt.format(mrev), size=9,
                               color=Palette.NAVY, font=Font.BODY, bold=True,
                               anchor=MSO_ANCHOR.MIDDLE)
                self._line(slide, MARGIN, Emu(y + row_h), tw,
                           color=Palette.GREY_LINE, weight=0.5)
                y = Emu(y + row_h)
        # total row
        self._line(slide, MARGIN, y, tw, color=Palette.NAVY, weight=1)
        if total is not None:
            self._text(slide, Emu(MARGIN + Inches(0.06)), y, Emu(bar_x0 - MARGIN),
                       row_h, total_label, size=9.5, color=Palette.NAVY,
                       font=Font.HEAD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            bw = bar_track  # the sum: a full-width reference bar
            self._rect(slide, bar_x0, Emu(y + row_h // 2 - Inches(0.08)), bw,
                       Inches(0.16), Palette.NAVY)
            self._text(slide, Emu(bar_x0 + bw + Inches(0.05)), y, Inches(1.0),
                       row_h, value_fmt.format(total), size=9.5, color=Palette.NAVY,
                       font=Font.HEAD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        return slide

    def save(self, path: str):
        self.prs.save(path)
        return path


__all__ = [
    "Deck", "Bullet", "ChartSpec", "MekkoColumn", "HarveyRow", "DriverNode",
    "DriverColumn", "MatrixItem", "RangeColumn", "RangeRow", "Palette", "Font",
]
